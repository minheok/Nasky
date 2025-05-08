import os
from pptx import Presentation
from docx import Document
import fitz  # PyMuPDF

def extract_text_from_pdf(file_path):
    text = ""
    with fitz.open(file_path) as doc:
        for page in doc:
            text += page.get_text()
    return text

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_file(file_path):
    if file_path.endswith(".pdf"):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith(".docx"):
        return extract_text_from_docx(file_path)
    elif file_path.endswith(".pptx"):
        return extract_text_from_pptx(file_path)
    else:
        return ""

def extract_texts_from_folder(folder_path):
    combined_text = ""
    for filename in os.listdir(folder_path):
        if filename.endswith((".pdf", ".docx", ".pptx")):
            file_path = os.path.join(folder_path, filename)
            file_text = extract_text_from_file(file_path)
            combined_text += f"\n\n[{filename}]\n{file_text}\n"
    return combined_text
